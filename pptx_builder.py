import argparse
import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt


def _set_widescreen(prs: Presentation) -> None:
    # 13.333" x 7.5" (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def _add_title(slide, text: str) -> None:
    # Small title at the top.
    left = Inches(0.6)
    top = Inches(0.2)
    width = Inches(12.2)
    height = Inches(0.6)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(18)


def _add_picture(prs: Presentation, slide, img_path: str) -> None:
    # Place the image below title.
    # Reduce size by ~30% vs previous width (12.2" -> 8.54"), and center it.
    top = Inches(1.0)
    width = Inches(12.2 * 0.7)
    # Center horizontally using presentation slide width.
    left = int((prs.slide_width - width) / 2)
    slide.shapes.add_picture(img_path, left, top, width=width)


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--manifest", required=True)
    ap.add_argument("--out", required=True)
    args = ap.parse_args()

    with open(args.manifest, "r", encoding="utf-8") as f:
        data = json.load(f)

    prs = Presentation()
    _set_widescreen(prs)

    slides = data.get("slides") or []
    for s in slides:
        title = (s.get("title") or "").strip()
        img = (s.get("image") or "").strip()
        if not img:
            continue

        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)
        if title:
            _add_title(slide, title)
        _add_picture(prs, slide, img)

    out_dir = os.path.dirname(os.path.abspath(args.out))
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    prs.save(args.out)
    print(f"OK: PPTX salvo em {args.out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
